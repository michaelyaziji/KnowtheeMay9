import streamlit as st
import os
from dotenv import load_dotenv
import tempfile
from document_processor import DocumentProcessor
from profile_generator import ProfileGenerator
from vector_store import VectorStore
from fpdf import FPDF
import re
import json
from pptx.dml.color import RGBColor
import pandas as pd
from pptx import Presentation
from io import BytesIO
import base64
from pathlib import Path
from employee_database import EmployeeDatabase
from query_processor import QueryProcessor
from rag_query_system import rag_system
import time

# Custom CSS for branding and layout
CUSTOM_CSS = """
<style>
/* Set initial zoom level and base font size */
html {
    zoom: 1.6;
    font-size: 20px;
}

body {
    background-color: #f8f9fa;
    font-size: 1.3rem; /* Larger base font size */
    color: #0a2c4d; /* Set default text color to blue */
}

/* Import Poppins font */
@import url('https://fonts.googleapis.com/css2?family=Poppins:wght@600&display=swap');

/* Additional custom font style */
@font-face {
    font-family: 'Poppins';
    font-style: normal;
    font-weight: 600;
    src: url(https://fonts.gstatic.com/s/poppins/v20/pxiByp8kv8JHgFVrLEj6Z1xlFd2JQEk.woff2) format('woff2');
}

.top-banner {
    position: fixed;
    left: 0;
    top: 0;
    width: 100vw;
    background: #0a2c4d;
    text-align: center;
    padding: 0.75rem;
    z-index: 100;
}

.banner-content {
    display: flex;
    justify-content: center;
    align-items: center;
}

.banner-text {
    color: #ffc72c !important;
    font-size: 3.2rem;
    font-weight: 600;
    font-family: 'Poppins', sans-serif !important;
    letter-spacing: 2.5px;
    text-align: center;
}

.top-spacer {
    height: 5.5rem;
}

.header-bar {
    background-color: #0a2c4d;
    padding: 2.5rem 2rem 2rem 2rem;
    color: white;
    border-radius: 0 0 24px 24px;
    margin-bottom: 2rem;
}

.header-title {
    font-size: 3.4rem;
    font-weight: 700;
    color: white;
    margin-bottom: 0.5rem;
}

.header-subtitle {
    font-size: 1.8rem;
    font-weight: 400;
    color: #e0e6ed;
    margin-bottom: 0.5rem;
}

.progress-cue {
    font-size: 1.6rem;
    color: #0a2c4d;
    margin-bottom: 2.2rem;
    font-weight: 500;
}

.section-title {
    color: #0a2c4d;
    font-size: 2.6rem;
    font-weight: 700;
    margin-top: 2.5rem;
    margin-bottom: 1rem;
}

.section-desc {
    color: #0a2c4d; /* Changed from #333 to blue */
    font-size: 1.6rem;
    margin-bottom: 0.7rem;
    margin-top: -0.5rem;
}

.profile-section {
    background: white;
    border-radius: 16px;
    padding: 2rem 2.5rem;
    margin-bottom: 2rem;
    box-shadow: 0 2px 8px rgba(10,44,77,0.07);
}

.footer {
    position: fixed;
    left: 0;
    bottom: 0;
    width: 100vw;
    background: #0a2c4d;
    color: #ffc72c;
    text-align: right;
    font-size: 1rem;
    font-weight: 600;
    padding: 0.23rem 2.5rem;
    letter-spacing: 1.5px;
    z-index: 100;
}

.bold-action {
    font-weight: 700;
    color: #0a2c4d;
}

.key-idea {
    color: #ffc72c;
    font-weight: 700;
}

/* Increase all Streamlit element font sizes */
.stCaption, .stMarkdown, .stTextInput, .stTextArea, .stFileUploader {
    font-size: 1.5rem !important;  /* Increased from 1.3rem */
    color: #0a2c4d !important; /* Set all text to blue */
}

/* Make buttons and form elements larger */
button, input, select, textarea, .stButton>button, .stSelectbox>div>div, [data-testid="stFileUploader"] {
    font-size: 1.5rem !important;
    color: #0a2c4d !important; /* Set all text to blue */
}

/* Make sure text in the file uploader is bigger */
[data-testid="stFileUploader"] span {
    font-size: 1.5rem !important;
    color: #0a2c4d !important; /* Set all text to blue */
}

/* Increase size of question text area */
textarea {
    min-height: 150px !important;  /* Increased from 120px */
    font-size: 1.5rem !important;
    color: #0a2c4d !important; /* Set all text to blue */
}

/* Make all text elements larger and blue */
p, div, span, label, a {
    font-size: 1.5rem;
    color: #0a2c4d !important; /* Set all text to blue */
}

/* Override any gray text */
[style*="color: #888"], [style*="color: rgb(136, 136, 136)"], 
[style*="color: gray"], [style*="color: #333"], [style*="color: rgb(51, 51, 51)"] {
    color: #0a2c4d !important; /* Override any inline gray styles */
}

/* Specific override for Streamlit text */
.css-nahz7x, .css-ffhzg2 {
    color: #0a2c4d !important;
}

/* Increase button size */
.stButton > button {
    padding: 0.8rem 1.5rem !important;
    height: auto !important;
    font-size: 1.5rem !important;
    min-height: 3rem !important;
}

/* Download buttons */
.stDownloadButton > button {
    padding: 0.8rem 1.5rem !important;
    font-size: 1.5rem !important;
    min-height: 3rem !important;
}

/* File uploader text */
.stFileUploader label {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
    font-weight: 600 !important;
}

.stFileUploader > div {
    font-size: 1.4rem !important;
}

/* Make expanders larger */
.streamlit-expanderHeader {
    font-size: 1.6rem !important;
    color: #0a2c4d !important; /* Set to blue */
    padding: 1rem !important;  /* Add more padding */
}

/* Streamlit tabs - make tab labels larger */
.stTabs [data-baseweb="tab-list"] button [data-testid="stMarkdownContainer"] p {
    font-size: 1.7rem !important;  /* Large tab text */
    font-weight: 600 !important;
    color: #0a2c4d !important;
}

/* Tab content area */
.stTabs [data-baseweb="tab-panel"] {
    font-size: 1.5rem !important;
}

/* Fix tab indicator positioning */
.stTabs [data-baseweb="tab-list"] {
    gap: 0 !important;
}

.stTabs [data-baseweb="tab-list"] button {
    border-radius: 0 !important;
    padding: 0.5rem 1rem !important;
    margin: 0 !important;
    position: relative !important;
}

/* Hide ALL default Streamlit tab indicators */
.stTabs [data-baseweb="tab-list"] button[aria-selected="true"] {
    border-bottom: none !important;
    border-bottom-color: transparent !important;
    box-shadow: none !important;
}

/* Hide the default indicator element if it exists */
.stTabs [data-baseweb="tab-list"] button::before {
    display: none !important;
}

/* Hide any other indicator elements */
.stTabs [data-baseweb="tab-list"] button > div::after,
.stTabs [data-baseweb="tab-list"] button > div::before {
    display: none !important;
}

/* Target the specific indicator that Streamlit creates */
.stTabs [data-baseweb="tab-list"] [role="tablist"] [role="tab"][aria-selected="true"] {
    border-bottom: none !important;
    border-bottom-color: transparent !important;
    border-bottom-width: 0 !important;
}

/* Hide any border or underline on tab content */
.stTabs [data-baseweb="tab-list"] [role="tab"] {
    border-bottom: none !important;
    text-decoration: none !important;
}

/* Active tab indicator - our custom one */
.stTabs [data-baseweb="tab-list"] button[aria-selected="true"]::after {
    content: "" !important;
    position: absolute !important;
    bottom: 0 !important;
    left: 0 !important;
    right: 0 !important;
    height: 3px !important;
    background-color: #ff4b4b !important;
    border-radius: 0 !important;
    z-index: 10 !important;
}

/* Ensure proper tab button styling */
.stTabs [data-baseweb="tab-list"] button {
    background: transparent !important;
    border: none !important;
    border-bottom: none !important;
}

/* Streamlit selectbox and other form elements */
.stSelectbox > div > div {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
}

.stSelectbox label {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
    font-weight: 600 !important;
}

/* Streamlit text input labels */
.stTextInput label, .stTextArea label {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
    font-weight: 600 !important;
}

/* Streamlit metrics */
[data-testid="metric-container"] {
    font-size: 1.4rem !important;
}

[data-testid="metric-container"] [data-testid="stMetricValue"] {
    font-size: 1.8rem !important;
    font-weight: 600 !important;
}

[data-testid="metric-container"] [data-testid="stMetricLabel"] {
    font-size: 1.3rem !important;
}

/* Streamlit dataframe */
.stDataFrame {
    font-size: 1.3rem !important;
}

.stDataFrame table {
    font-size: 1.3rem !important;
}

.stDataFrame th {
    font-size: 1.4rem !important;
    font-weight: 600 !important;
}

/* Make markdown headers larger */
.main h1 {
    font-size: 2.8rem !important;
    color: #0a2c4d !important;
}

.main h2 {
    font-size: 2.4rem !important;
    color: #0a2c4d !important;
}

.main h3 {
    font-size: 2.0rem !important;
    color: #0a2c4d !important;
}

.main h4 {
    font-size: 1.8rem !important;
    color: #0a2c4d !important;
}

/* Additional specific targeting */
[data-testid="stMarkdownContainer"] p {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
}

[data-testid="stMarkdownContainer"] li {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
}

[data-testid="stMarkdownContainer"] strong {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
}

/* Enhanced Loading animation */
@keyframes bouncing {
    0%, 20%, 50%, 80%, 100% { transform: translateY(0); }
    40% { transform: translateY(-8px); }
    60% { transform: translateY(-4px); }
}

@keyframes pulsing {
    0% { opacity: 0.7; }
    50% { opacity: 1; }
    100% { opacity: 0.7; }
}

.loading-indicator {
    position: fixed;
    left: 50%;
    top: 50%;
    transform: translate(-50%, -50%);
    background: linear-gradient(135deg, #0a2c4d 0%, #1a4a7a 100%);
    color: #ffc72c !important;
    padding: 16px 24px;
    border-radius: 25px;
    z-index: 10000;
    font-weight: bold;
    font-size: 16px;
    box-shadow: 0 8px 32px rgba(10, 44, 77, 0.4);
    display: flex;
    align-items: center;
    border: 2px solid #ffc72c;
    animation: pulsing 1.5s infinite;
    backdrop-filter: blur(5px);
    min-width: 200px;
    text-align: center;
    justify-content: center;
}

.loading-indicator::before {
    content: "🏃‍♂️";
    font-size: 22px;
    margin-right: 10px;
    animation: bouncing 1.2s infinite;
    display: inline-block;
}

.loading-indicator.processing::before {
    content: "⚡";
    animation: bouncing 0.8s infinite;
}

#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}

/* Streamlit info, success, warning, error messages */
.stAlert {
    font-size: 1.4rem !important;
}

.stSuccess, .stInfo, .stWarning, .stError {
    font-size: 1.4rem !important;
}

/* Spinner text */
.stSpinner > div {
    font-size: 1.4rem !important;
}

/* Checkbox and radio button labels */
.stCheckbox label, .stRadio label {
    font-size: 1.5rem !important;
    color: #0a2c4d !important;
}

/* Expander content */
.streamlit-expanderContent {
    font-size: 1.5rem !important;
    padding: 1rem !important;
}

/* Search results and employee cards */
.stExpander {
    font-size: 1.5rem !important;
}

.stExpander > div > div {
    font-size: 1.5rem !important;
}

/* Force larger fonts on specific Streamlit containers */
.main .block-container {
    font-size: 1.5rem !important;
}

/* Streamlit columns content */
.block-container .element-container {
    font-size: 1.5rem !important;
}
</style>

<script>
// Enhanced loading indicator management
let loadingIndicator;
let loadingTimeout;
let isProcessing = false;

function showLoadingIndicator(message = 'Processing...', type = 'default') {
    if (!loadingIndicator) return;
    
    clearTimeout(loadingTimeout);
    isProcessing = true;
    
    loadingIndicator.textContent = message;
    loadingIndicator.className = `loading-indicator ${type}`;
    loadingIndicator.style.display = 'flex';
    
    console.log('🔄 Loading indicator shown:', message);
}

function hideLoadingIndicator(delay = 500) {
    if (!loadingIndicator) return;
    
    clearTimeout(loadingTimeout);
    loadingTimeout = setTimeout(() => {
        if (loadingIndicator && !document.querySelector('.stSpinner')) {
            loadingIndicator.style.display = 'none';
            isProcessing = false;
            console.log('✅ Loading indicator hidden');
        }
    }, delay);
}

// Set initial zoom level using JavaScript as a backup
document.body.style.zoom = "160%";

// Enhanced loading indicator initialization
document.addEventListener('DOMContentLoaded', function() {
    // Find all elements with gray text and change to blue
    const grayTexts = document.querySelectorAll('[style*="color: #888"], [style*="color: rgb(136, 136, 136)"], [style*="color: gray"], [style*="color: #333"], [style*="color: rgb(51, 51, 51)"]');
    grayTexts.forEach(element => {
        element.style.color = '#0a2c4d';
    });
    
    // Create enhanced loading indicator
    loadingIndicator = document.createElement('div');
    loadingIndicator.className = 'loading-indicator';
    loadingIndicator.textContent = 'Processing...';
    loadingIndicator.style.display = 'none';
    document.body.appendChild(loadingIndicator);
    
    // Monitor all button clicks
    document.addEventListener('click', function(e) {
        const button = e.target.closest('button');
        if (button) {
            const buttonText = button.textContent.trim();
            console.log('🖱️ Button clicked:', buttonText);
            
            // Determine message based on button
            let message = 'Processing...';
            if (buttonText.includes('Generate')) {
                message = 'Generating Profile...';
            } else if (buttonText.includes('Submit')) {
                message = 'Analyzing Question...';
            } else if (buttonText.includes('Download')) {
                message = 'Preparing Download...';
            } else if (buttonText.includes('Import')) {
                message = 'Importing Data...';
            } else if (buttonText.includes('Search')) {
                message = 'Searching Database...';
            }
            
            showLoadingIndicator(message, 'processing');
        }
    });
    
    // Monitor file uploads
    document.addEventListener('change', function(e) {
        if (e.target.type === 'file' && e.target.files.length > 0) {
            console.log('📁 Files uploaded:', e.target.files.length);
            showLoadingIndicator('Processing Documents...', 'processing');
        }
    });
    
    // Enhanced mutation observer for Streamlit elements
    const observer = new MutationObserver(function(mutations) {
        const spinner = document.querySelector('.stSpinner');
        const progress = document.querySelector('.stProgress');
        const status = document.querySelector('[data-testid="stStatus"]');
        const alertElements = document.querySelectorAll('[data-testid="stAlert"]');
        
        if (spinner || progress || status) {
            if (!isProcessing) {
                showLoadingIndicator('Working...', 'processing');
            }
        } else if (isProcessing) {
            // Check if any error/success alerts appeared
            let hasAlert = false;
            alertElements.forEach(alert => {
                const alertText = alert.textContent.toLowerCase();
                if (alertText.includes('success') || alertText.includes('error') || 
                    alertText.includes('complete') || alertText.includes('failed')) {
                    hasAlert = true;
                }
            });
            
            if (hasAlert) {
                hideLoadingIndicator(200); // Hide quickly if there's a result
            } else {
                hideLoadingIndicator();
            }
        }
    });
    
    observer.observe(document.body, { 
        childList: true, 
        subtree: true,
        attributes: true,
        attributeFilter: ['class', 'style']
    });
    
    // Initial page load
    showLoadingIndicator('Loading Application...', 'default');
    setTimeout(() => {
        if (!document.querySelector('.stSpinner')) {
            hideLoadingIndicator(1000);
        }
    }, 2000);
    
    console.log('🚀 Enhanced loading indicator initialized');
});

// Enhanced download function
window.downloadReference = function(fileName, employeeId) {
    showLoadingIndicator('Preparing Download...', 'processing');
    
    sessionStorage.setItem('downloadReference', JSON.stringify({
        fileName: fileName,
        employeeId: employeeId,
        timestamp: Date.now()
    }));
    
    const dummyButton = document.createElement('button');
    dummyButton.style.display = 'none';
    dummyButton.id = 'download-trigger-' + Date.now();
    document.body.appendChild(dummyButton);
    dummyButton.click();
    document.body.removeChild(dummyButton);
};
</script>
"""

# Load environment variables
load_dotenv()

# Initialize session state
for key, default in {
    'subject_docs': [],
    'context_docs': [],
    'team_docs': [],
    'profile': None,
    'user_question': '',
    'question_answer': None,
    'reference_docs': [],
    'developer_mode': False,
    'intent': "Get an overall assessment",
    'intent_other': '',
    'document_cache': {},  # Cache for document bytes to enable reference downloads
    'document_download_request': None  # To track download requests from JavaScript
}.items():
    if key not in st.session_state:
        st.session_state[key] = default

# Initialize components
document_processor = DocumentProcessor()
profile_generator = ProfileGenerator()

# Load and process reference PDFs from HowToInterpret/
def load_reference_docs():
    reference_folder = "HowToInterpret"
    reference_texts = []
    for filename in os.listdir(reference_folder):
        if filename.lower().endswith('.pdf'):
            file_path = os.path.join(reference_folder, filename)
            try:
                text, metadata = document_processor.process_document(file_path)
                reference_texts.append(text)
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
    return reference_texts

# Only load reference docs once per session
if not st.session_state.reference_docs:
    st.session_state.reference_docs = load_reference_docs()

def create_pdf(profile_text, question_answer=None):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    pdf.add_page()

    # Try to use a custom font, but fall back to a standard font if not available
    try:
        # Robust font path resolution
        BASE_DIR = Path(__file__).resolve().parent
        FONT_PATH = BASE_DIR / "fonts" / "DejaVuSans.ttf"
        
        if FONT_PATH.exists():
            pdf.add_font("DejaVu", "", str(FONT_PATH), uni=True)
            pdf.set_font("DejaVu", size=12)
        else:
            # Fall back to standard font
            pdf.set_font("Helvetica", size=12)
    except Exception as e:
        print(f"Font loading error: {e}")
        # Use standard font as fallback
        pdf.set_font("Helvetica", size=12)

    lines = profile_text.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            continue  # Skip empty lines
        # Break up very long words (over 80 chars)
        if max((len(word) for word in line.split()), default=0) > 80:
            words = [word if len(word) <= 80 else ' '.join([word[i:i+80] for i in range(0, len(word), 80)]) for word in line.split()]
            line = ' '.join(words)
        try:
            if line.startswith("###"):
                heading = line.replace("###", "").strip()
                pdf.set_font_size(14)
                pdf.set_font("", "B")  # Bold
                pdf.cell(0, 10, heading, ln=True)
                pdf.set_font_size(12)
                pdf.set_font("", "")  # Normal
            elif re.match(r"^\s*\d+\.", line) or line.startswith("-"):
                pdf.cell(10)
                pdf.multi_cell(0, 8, line)
            elif "**" in line:
                parts = re.split(r'(\*\*.*?\*\*)', line)
                for part in parts:
                    if part.startswith("**") and part.endswith("**"):
                        pdf.set_font("", "B")  # Bold
                        pdf.write(8, part[2:-2])
                        pdf.set_font("", "")  # Normal
                    else:
                        pdf.write(8, part)
                pdf.ln(10)
            else:
                pdf.multi_cell(0, 8, line)
        except Exception as e:
            print(f"Skipped line in PDF due to error: {e}\nLine content: {repr(line)}")
            continue


    if question_answer:
        pdf.ln(10)
        pdf.set_font("", "B")  # Bold
        pdf.cell(0, 10, "Special Question Answer:", ln=True)
        pdf.set_font("", "")  # Normal

        try:
            pdf.multi_cell(0, 10, question_answer)
        except Exception as e:
            print(f"Skipped question_answer in PDF due to error: {e}")

    return pdf.output(dest='S')

def clean_source_text(source_text):
    """Clean up temporary filenames in sources text and replace generic file types with meaningful document descriptions."""
    if not source_text:
        return ""
    
    # First, extract source types from the text if possible
    source_types = []
    
    # Check for Hogan references
    if "Hogan" in source_text or "hogan" in source_text:
        source_types.append("Hogan Assessment")
    
    # Check for IDI references - distinguish between the two IDI types
    if "IDI" in source_text or "idi" in source_text:
        # Check for Individual Directions Inventory
        if "directions" in source_text.lower() or "individual directions" in source_text.lower():
            source_types.append("Individual Directions Inventory")
        # Check for Intercultural Development Assessment
        elif "intercultural" in source_text.lower() or "cultural" in source_text.lower():
            source_types.append("Intercultural Development Assessment")
        # Generic IDI if we can't determine which one
        else:
            source_types.append("Assessment")
    
    # Check for 360 references
    if "360" in source_text:
        source_types.append("360° Feedback")
    
    # Check for CV references
    if "CV" in source_text or "cv" in source_text or "resume" in source_text.lower():
        source_types.append("CV/Resume")
    
    # Replace specific patterns
    # Replace Hogan temp files
    source_text = re.sub(r'tmp[a-zA-Z0-9]+\.pdf\s*\(Hogan\)', 'Hogan Assessment', source_text)
    
    # Replace Individual Directions Inventory files
    source_text = re.sub(r'tmp[a-zA-Z0-9]+\.pdf\s*\(IDI\)', 'Individual Directions Inventory', source_text)
    source_text = re.sub(r'tmp[a-zA-Z0-9]+\.pdf\s*\(Individual Directions\)', 'Individual Directions Inventory', source_text)
    
    # Replace Intercultural Development references
    source_text = re.sub(r'tmp[a-zA-Z0-9]+\.pdf\s*\(Intercultural\)', 'Intercultural Development Assessment', source_text)
    
    # Replace other temp files
    source_text = re.sub(r'tmp[a-zA-Z0-9]+\.[a-z]+', '', source_text)
    source_text = re.sub(r'tmp[a-zA-Z0-9]+', '', source_text)
    
    # Replace generic file types with more meaningful descriptions
    source_text = re.sub(r'\bPDF\b', 'Document', source_text)
    source_text = re.sub(r'\bDOCX\b', 'Document', source_text)
    source_text = re.sub(r'\bDOC\b', 'Document', source_text)
    
    # Clean up formatting
    source_text = re.sub(r'\s+', ' ', source_text)  # Multiple spaces
    source_text = re.sub(r',\s*,', ',', source_text)  # Multiple commas
    source_text = re.sub(r'\(\s*\)', '', source_text)  # Empty parentheses
    source_text = re.sub(r',\s*$', '', source_text)  # Trailing commas
    source_text = re.sub(r'^\s*,\s*', '', source_text)  # Leading commas
    source_text = re.sub(r'\(\s*,', '(', source_text)  # Commas after opening parenthesis
    source_text = re.sub(r',\s*\)', ')', source_text)  # Commas before closing parenthesis
    
    # Final cleanup
    source_text = source_text.strip()
    
    # If we've removed everything but have identified source types, use them
    if (not source_text or source_text == ',' or source_text == '()') and source_types:
        return ", ".join(source_types)
    
    # If we still have nothing, return a generic source
    if not source_text or source_text == ',' or source_text == '()':
        return "Assessment Documents"
    
    return source_text

def generate_pptx_from_json(json_data, template_path=None):
    """
    Generate a PowerPoint presentation from structured JSON data.
    Uses a template if provided, otherwise creates a new presentation.
    Maps each section to the appropriate slide in the template.
    """
    # Load template if provided, otherwise use blank
    if template_path:
        try:
            prs = Presentation(template_path)
            print(f"Using template: {template_path}")
            # Debug template info
            print(f"Template has {len(prs.slides)} slides")
            for i, slide in enumerate(prs.slides):
                print(f"Slide {i+1}: {len(slide.shapes)} shapes")
                
            # Debug theme colors
            if hasattr(prs, 'theme') and hasattr(prs.theme, 'theme_color_scheme'):
                print("\nTHEME COLORS:")
                for idx, color in enumerate(prs.theme.theme_color_scheme.colors):
                    if hasattr(color, 'rgb'):
                        r, g, b = color.rgb >> 16, (color.rgb >> 8) & 0xFF, color.rgb & 0xFF
                        print(f"  Color {idx+1}: RGB({r},{g},{b})")
                    else:
                        print(f"  Color {idx+1}: [No RGB value available]")
            else:
                print("\nNo theme color information available")
        except Exception as e:
            print(f"Error loading template: {e}")
            prs = Presentation()
    else:
        prs = Presentation()
        
    # Define our brand colors
    HEADER_COLOR_WHITE = RGBColor(255, 255, 255)  # White for headers
    BODY_COLOR_BLUE = RGBColor(10, 44, 77)        # Deep blue for body text - matches template
    
    # If using a blank presentation, create slides for each section
    if template_path is None or len(prs.slides) < 2:  # If no template or not enough slides
        for section in json_data:
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            # Set title if possible
            title_shape = None
            for shape in slide.shapes:
                if shape.name.startswith('Title'):
                    title_shape = shape
                    break
            if title_shape and hasattr(title_shape, 'text_frame'):
                title_shape.text_frame.text = section['section']
                # Set title font to white for contrast against blue background
                for paragraph in title_shape.text_frame.paragraphs:
                    if hasattr(paragraph.font, 'color') and hasattr(paragraph.font.color, 'rgb'):
                        paragraph.font.color.rgb = HEADER_COLOR_WHITE
                    if hasattr(paragraph.font, 'bold'):
                        paragraph.font.bold = True
                    if hasattr(paragraph.font, 'size'):
                        paragraph.font.size = 32 * 12700  # 32pt for headers
            
            # Add content
            content = section['content']
            sources = section.get('sources', '')
            
            # Clean up sources to remove temporary filenames
            sources = clean_source_text(sources)
            
            content_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # MSO_PLACEHOLDER.BODY
                    content_shape = shape
                    break
            if content_shape:
                content_shape.text_frame.clear()
                for line in content.split('\n'):
                    if not line.strip():
                        continue  # skip empty lines
                    p = content_shape.text_frame.add_paragraph()
                    p.text = line.lstrip('-* ').strip()
                    if line.strip().startswith(('-', '*')):
                        p.level = 0
                        p.font.bullet = True
                    if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                        try:
                            p.font.color.rgb = RGBColor(10, 44, 77)
                        except:
                            pass
                    if hasattr(p.font, 'name'):
                        try:
                            p.font.name = "Calibri"
                        except:
                            pass
                    if hasattr(p.font, 'size'):
                        try:
                            p.font.size = 18 * 12700  # 18pt (increased from 12pt)
                        except:
                            pass
                # Add sources as a separate paragraph if present
                if sources:
                    p = content_shape.text_frame.add_paragraph()
                    p.text = f"Sources: {sources}"
                    if hasattr(p.font, 'italic'):
                        p.font.italic = True
                    if hasattr(p.font, 'name'):
                        p.font.name = "Calibri"
                    if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                        p.font.color.rgb = RGBColor(10, 44, 77)
                    if hasattr(p.font, 'size'):
                        p.font.size = 18 * 12700  # 18pt (increased from 12pt)
    else:
        # Using the template - map sections to specific slides
        # Based on the template structure shown in screenshots
        
        # Keep cover slide (slide 1) and table of contents (slide 2) as is
        
        # Map each section to the corresponding template slide
        # Slide indexing starts at 0
        section_map = {
            "Profile Summary": 2,         # Slide 3
            "Leadership Summary": 2,      # Slide 3 (alternate name)
            "Key Strengths": 3,           # Slide 4
            "Potential Derailers": 4,     # Slide 5 
            "Leadership Style": 5,        # Slide 6
            "Role Fit Chart - Good Fit": 6,  # Slide 7
            "Role Fit Chart - Poor Fit": 7,   # Slide 8
            "Role Fit Chart - Bad Fit": 7,    # Slide 8 (using "Bad" instead of "Poor")
            "Role-fit Chart: Good fit": 6,   # Alternate spelling
            "Role-fit Chart: Poor fit": 7,   # Alternate spelling
            "Roles That Would Fit": 6,    # Alternate wording (slide 7)
            "Roles That Would Not Fit": 7, # Alternate wording (slide 8)
            "Good Role Fit": 6,           # Another variation (slide 7)
            "Poor Role Fit": 7,           # Another variation (slide 8)
            "Bad Role Fit": 7,            # Another variation (slide 8)
            "Role Fit: Good": 6,          # Another variation (slide 7)
            "Role Fit: Poor": 7,          # Another variation (slide 8)
            "Role Fit: Bad": 7,           # Another variation (slide 8)
            "Good Fit Roles": 6,          # Another variation (slide 7)
            "Bad Fit Roles": 7,           # Another variation (slide 8)
            "Poor Fit Roles": 7,          # Another variation (slide 8)
            "Special Queries, if any": 8,    # Slide 9
            "Special Query": 8,             # Alternate name
        }
        
        # Process each section from the JSON data
        print(f"Processing {len(json_data)} sections")
        for section in json_data:
            section_name = section['section']
            print(f"Processing section: '{section_name}'")
                
            content = section['content']
            sources = section.get('sources', '')
            
            # Clean up sources to remove temporary filenames
            sources = clean_source_text(sources)
            
            # Find the slide index for this section
            slide_idx = section_map.get(section_name)
            if slide_idx is None:
                # Try more flexible matching for alternate section names
                for map_name, idx in section_map.items():
                    # Check for strict containment first
                    if map_name.lower() in section_name.lower() or section_name.lower() in map_name.lower():
                        slide_idx = idx
                        print(f"Found slide match: '{section_name}' -> '{map_name}' (slide {idx+1})")
                        break
                
                # If still no match, try keyword matching for role fit sections
                if slide_idx is None:
                    section_lower = section_name.lower()
                    # Check for good/positive fit keywords
                    if any(word in section_lower for word in ['good fit', 'would fit', 'positive', 'strong match']):
                        slide_idx = 6  # Slide 7 - Good fit
                        print(f"Found keyword match for good fit: '{section_name}' (slide 7)")
                    # Check for poor/bad fit keywords
                    elif any(word in section_lower for word in ['poor fit', 'bad fit', 'not fit', 'wouldn\'t fit', 'negative']):
                        slide_idx = 7  # Slide 8 - Poor fit
                        print(f"Found keyword match for poor fit: '{section_name}' (slide 8)")
                    # Special query match
                    elif any(word in section_lower for word in ['special', 'query', 'question', 'answer']):
                        slide_idx = 8  # Slide 9 - Special query
                        print(f"Found keyword match for special query: '{section_name}' (slide 9)")
            
            if slide_idx is not None and slide_idx < len(prs.slides):
                # Use existing slide from template
                slide = prs.slides[slide_idx]
                print(f"Adding content to slide {slide_idx+1} for section '{section_name}'")
                
                # IMPORTANT CHANGE: SKIP ALL TITLE MANIPULATION
                # We will leave the template titles exactly as they are
                
                # FIND OR CREATE CONTENT SHAPE
                content_shape = None
                
                # Look for existing content shapes (not the title)
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        # Skip any shape that looks like a title (usually at the top of slide)
                        if hasattr(shape, 'top') and shape.top < 1000000:  # ~1 inch from top
                            continue
                        # Use the first non-title text shape we find for content
                        content_shape = shape
                        print(f"Found content shape in slide {slide_idx+1}")
                        break
                
                # If no content shape found, create a new textbox for content
                if not content_shape:
                    try:
                        # Create new textbox with better positioning (below title)
                        content_left = 0.5 * 914400    # 0.5 inch from left
                        content_top = 650000     # 0.5 inches from top (below title)
                        content_width = prs.slide_width - (1 * 914400)  # Full width minus 1 inch
                        content_height = prs.slide_height - content_top - (0.5 * 914400)  # From top to bottom with margin
                        
                        content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                        print(f"Created new content textbox on slide {slide_idx+1}")
                    except Exception as e:
                        print(f"Error creating content textbox: {e}")
                
                if content_shape:
                    try:
                        # Clear existing text
                        if hasattr(content_shape, 'text_frame'):
                            content_shape.text_frame.clear()
                        
                        # Adjust position for better spacing from header
                        if hasattr(content_shape, 'top'):
                            try:
                                if content_shape.top < 650000:
                                    content_shape.top = 650000
                            except:
                                pass
                        
                        # Enable word wrap 
                        if hasattr(content_shape, 'text_frame') and hasattr(content_shape.text_frame, 'word_wrap'):
                            content_shape.text_frame.word_wrap = True
                        
                        # Add margins to text frame if possible
                        if hasattr(content_shape, 'text_frame') and hasattr(content_shape.text_frame, 'margin_top'):
                            try:
                                # Add generous margins
                                content_shape.text_frame.margin_top = 300000     # ~8mm top margin
                                content_shape.text_frame.margin_bottom = 150000  # ~4mm bottom margin
                                content_shape.text_frame.margin_left = 150000    # ~4mm left margin
                                content_shape.text_frame.margin_right = 150000   # ~4mm right margin
                            except:
                                pass
                        
                        # Add main content as bullet-aware paragraphs
                        for line in content.split('\n'):
                            if not line.strip():
                                continue
                            p = content_shape.text_frame.add_paragraph()
                            p.text = line.lstrip('-* ').strip()
                            if line.strip().startswith(('-', '*')):
                                p.level = 0
                                p.font.bullet = True
                            if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                                try:
                                    p.font.color.rgb = RGBColor(10, 44, 77)
                                except:
                                    pass
                            if hasattr(p.font, 'name'):
                                try:
                                    p.font.name = "Calibri"
                                except:
                                    pass
                            if hasattr(p.font, 'size'):
                                try:
                                    p.font.size = 18 * 12700  # 18pt (increased from 12pt)
                                except:
                                    pass
                        
                        # Add sources as a separate paragraph if present
                        if sources:
                            p = content_shape.text_frame.add_paragraph()
                            p.text = f"Sources: {sources}"
                            if hasattr(p.font, 'italic'):
                                p.font.italic = True
                            if hasattr(p.font, 'name'):
                                p.font.name = "Calibri"
                            if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                                p.font.color.rgb = RGBColor(10, 44, 77)
                            if hasattr(p.font, 'size'):
                                p.font.size = 18 * 12700  # 18pt (increased from 12pt)
                        
                        print(f"Successfully added content to slide {slide_idx+1}")
                    except Exception as e:
                        print(f"Error setting content on slide {slide_idx+1}: {e}")
            else:
                print(f"No matching slide found for section '{section_name}', creating new slide")
                # Add a new slide with appropriate formatting
                slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                
                # Create proper title with blue background
                try:
                    # Create a title box with blue background
                    title_left = 0
                    title_top = 0
                    title_width = prs.slide_width
                    title_height = 0.8 * 914400  # 0.8 inches
                    
                    # Add shape for title background
                    title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                    
                    # Apply blue background
                    if hasattr(title_shape, 'fill'):
                        title_shape.fill.solid()
                        title_shape.fill.fore_color.rgb = RGBColor(10, 44, 77)
                    
                    # Add title text
                    title_shape.text_frame.clear()
                    p = title_shape.text_frame.add_paragraph()
                    p.text = section_name
                    p.alignment = 1  # Center
                    
                    # Format title - white text
                    if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                        p.font.color.rgb = RGBColor(255,255,255)
                    if hasattr(p.font, 'bold'):
                        p.font.bold = True
                    if hasattr(p.font, 'size'):
                        p.font.size = 36 * 12700  # 36pt (increased from 32pt)
                        
                    # Apply to all runs
                    for run in p.runs:
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                            run.font.color.rgb = RGBColor(255,255,255)
                        if hasattr(run.font, 'bold'):
                            run.font.bold = True
                        if hasattr(run.font, 'size'):
                            run.font.size = 36 * 12700  # 36pt (increased from 32pt)
                except Exception as e:
                    print(f"Error creating title on new slide: {e}")
                
                # Create content box
                try:
                    # Create content box
                    content_left = 0.5 * 914400    # 0.5 inch from left
                    content_top = 650000     # 1.2 inches from top (below title)
                    content_width = prs.slide_width - (1 * 914400)  # Full width minus 1 inch margins
                    content_height = prs.slide_height - content_top - (0.5 * 914400)  # To bottom with margin
                    
                    content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                    
                    # Add margins to text frame
                    if hasattr(content_shape, 'text_frame') and hasattr(content_shape.text_frame, 'margin_top'):
                        content_shape.text_frame.margin_top = 300000     # Top margin
                        content_shape.text_frame.margin_bottom = 150000  # Bottom margin
                        content_shape.text_frame.margin_left = 150000    # Left margin
                        content_shape.text_frame.margin_right = 150000   # Right margin
                    
                    # Enable word wrap
                    if hasattr(content_shape, 'text_frame') and hasattr(content_shape.text_frame, 'word_wrap'):
                        content_shape.text_frame.word_wrap = True
                    
                    # Add content
                    p = content_shape.text_frame.add_paragraph()
                    p.text = content
                    
                    # Format content - blue text
                    if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                        p.font.color.rgb = RGBColor(10, 44, 77)
                    if hasattr(p.font, 'name'):
                        p.font.name = "Calibri"
                    if hasattr(p.font, 'size'):
                        p.font.size = 18 * 12700  # 18pt (increased from 12pt)
                        
                    # Apply to all runs
                    for run in p.runs:
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                            run.font.color.rgb = RGBColor(10, 44, 77)
                        if hasattr(run.font, 'name'):
                            run.font.name = "Calibri"
                        if hasattr(run.font, 'size'):
                            run.font.size = 18 * 12700  # 18pt (increased from 12pt)
                    
                    # Add sources if available
                    if sources:
                        p = content_shape.text_frame.add_paragraph()
                        p.text = f"Sources: {sources}"
                        
                        # Format sources - italic
                        if hasattr(p.font, 'italic'):
                            p.font.italic = True
                        if hasattr(p.font, 'name'):
                            p.font.name = "Calibri"
                        if hasattr(p.font, 'color') and hasattr(p.font.color, 'rgb'):
                            p.font.color.rgb = RGBColor(10, 44, 77)
                        if hasattr(p.font, 'size'):
                            p.font.size = 18 * 12700  # 18pt (increased from 12pt)
                except Exception as e:
                    print(f"Error creating content on new slide: {e}")
    
    # Save to BytesIO
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def display_profile(profile_json):
    """Display the leadership profile in a structured format."""
    try:
        # Parse the profile JSON
        profile_data = json.loads(profile_json)
        
        st.markdown('<div class="section-title">Leadership Profile</div>', unsafe_allow_html=True)
        
        # Display each section of the profile
        for section in profile_data:
            section_name = section.get("section", "")
            content = section.get("content", "")
            sources = section.get("sources", "")
            
            st.markdown(f'<div class="profile-section">', unsafe_allow_html=True)
            st.markdown(f"### {section_name}")
            st.markdown(content)
            
            if sources:
                st.markdown(f"<small><i>Sources: {sources}</i></small>", unsafe_allow_html=True)
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    except Exception as e:
        st.error(f"Error displaying profile: {e}")
        st.markdown('<div class="profile-section">', unsafe_allow_html=True)
        st.markdown(profile_json)
        st.markdown('</div>', unsafe_allow_html=True)

def generate_pptx_from_profile(profile_json):
    """Generate a PowerPoint presentation from the profile JSON."""
    # Parse the profile JSON if it's a string
    if isinstance(profile_json, str):
        profile_data = json.loads(profile_json)
    else:
        profile_data = profile_json
    
    # Get the template path
    template_path = Path(__file__).resolve().parent / "template.pptx"
    if not template_path.exists():
        template_path = None
    
    # Generate PowerPoint
    return generate_pptx_from_json(profile_data, template_path=template_path)

def create_reference_links(employee_data, reference_texts):
    """
    Create clickable download links for reference documents mentioned in profile answers.
    
    Args:
        employee_data: Dictionary containing employee metadata and profile
        reference_texts: List of reference names mentioned in the answer
        
    Returns:
        HTML string with clickable download links for each reference
    """
    # Common reference types and their possible variations
    reference_mapping = {
        "CV": ["CV", "CV/Resume", "Resume", "Curriculum Vitae", "resume"],
        "Hogan": ["Hogan", "Hogan Assessment", "HPI", "HDS", "MVPI"],
        "IDI": ["IDI", "IDI Assessment", "Individual Directions Inventory", "Intercultural Development Assessment"],
        "NEO": ["NEO", "NEO Assessment", "NEO-PI", "NEO-PI-R", "Big Five"],
        "360": ["360", "360 Assessment", "360° Feedback", "360 Feedback"]
    }
    
    # Create HTML for links
    links_html = ""
    
    # If we have document names stored, use them
    if 'document_names' in employee_data['metadata'] and employee_data['metadata']['document_names']:
        document_names = employee_data['metadata']['document_names']
        
        # Create a mapping of reference type to actual filename
        file_mapping = {}
        for doc_name in document_names:
            doc_lower = doc_name.lower()
            for ref_type, variations in reference_mapping.items():
                if any(var.lower() in doc_lower for var in variations):
                    file_mapping[ref_type] = doc_name
                    break
        
        # Create a link for each reference text based on the mapping
        for ref_text in reference_texts:
            ref_text = ref_text.strip()
            file_name = None
            
            # Try to match this reference to a file
            for ref_type, variations in reference_mapping.items():
                if any(var.lower() in ref_text.lower() for var in variations):
                    if ref_type in file_mapping:
                        file_name = file_mapping[ref_type]
                        break
            
            # Create a link if we found a matching file
            if file_name:
                links_html += f'<li><a href="#" onclick="downloadReference(\'{file_name}\', \'{employee_data["id"]}\');">{ref_text}</a></li>\n'
            else:
                links_html += f'<li>{ref_text}</li>\n'
    else:
        # If no document names are stored, just display plain text
        for ref_text in reference_texts:
            links_html += f'<li>{ref_text}</li>\n'
    
    return f'<ul>\n{links_html}</ul>\n'

def cache_document(doc, employee_id=None):
    """
    Cache a document in session state for later download as a reference.
    
    Args:
        doc: A file-like object with a 'name' attribute and getvalue() method
        employee_id: Optional employee ID to associate with the document
    """
    # Create a unique key for this document
    doc_key = f"{employee_id}_{doc.name}" if employee_id else doc.name
    
    # Store document bytes in cache
    st.session_state.document_cache[doc_key] = {
        'bytes': doc.getvalue(),
        'name': doc.name,
        'employee_id': employee_id,
        'mime_type': get_mime_type(doc.name)
    }
    
def get_mime_type(filename):
    """Determine MIME type based on file extension."""
    ext = filename.lower().split('.')[-1]
    mime_types = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'doc': 'application/msword',
        'txt': 'text/plain',
        'csv': 'text/csv',
        'xls': 'application/vnd.ms-excel',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    }
    return mime_types.get(ext, 'application/octet-stream')

def handle_document_download_request():
    """
    Check for and handle document download requests from JavaScript.
    This function should be called at the beginning of the main UI function.
    """
    # Create a container for download buttons
    download_container = st.empty()
    
    # Check for query parameters (Streamlit doesn't support directly reading JS sessionStorage)
    # Instead we rely on a hidden button click from JS to trigger a rerun
    
    # Use JavaScript to inject query parameter values
    st.markdown("""
    <script>
    // Check for document download request in sessionStorage
    document.addEventListener('DOMContentLoaded', function() {
        const downloadRequest = sessionStorage.getItem('downloadReference');
        if (downloadRequest) {
            // Parse download request
            const request = JSON.parse(downloadRequest);
            
            // Set custom data attribute on body for server to read
            document.body.setAttribute('data-download-request', JSON.stringify(request));
            
            // Create a mutation observer to watch for the download button
            const observer = new MutationObserver(function(mutations) {
                // Look for the added download button
                const downloadButton = document.querySelector('[data-testid="stDownloadButton"]');
                if (downloadButton) {
                    // Click it automatically
                    downloadButton.click();
                    
                    // Clean up
                    sessionStorage.removeItem('downloadReference');
                    document.body.removeAttribute('data-download-request');
                    observer.disconnect();
                }
            });
            
            // Start observing
            observer.observe(document.body, { childList: true, subtree: true });
        }
    });
    </script>
    """, unsafe_allow_html=True)
    
    # Use session state to track if we've handled a download request
    if 'doc_download_handled' not in st.session_state:
        st.session_state.doc_download_handled = False
    
    # Try to extract the download request from custom HTML attributes
    # This is a workaround since we can't directly read JavaScript sessionStorage
    
    # Execute JavaScript to check for download request
    js_code = """
    if (document.body.hasAttribute('data-download-request')) {
        return document.body.getAttribute('data-download-request');
    }
    return '';
    """
    
    try:
        # We can't actually run this code here, but if we could:
        # download_request_json = st.experimental_get_query_params().get('download_request', [''])[0]
        # For now, we'll check if there's an uploaded file in the cache that matches by name
        
        # For each document in the cache, create a hidden download button
        if not st.session_state.doc_download_handled and st.session_state.document_cache:
            # Only try to download the most recent uploaded document as a fallback
            # In a real implementation, we'd match the exact requested document
            latest_doc_key = list(st.session_state.document_cache.keys())[-1]
            doc_info = st.session_state.document_cache[latest_doc_key]
            
            with download_container:
                st.download_button(
                    label=f"Download {doc_info['name']}",
                    data=doc_info['bytes'],
                    file_name=doc_info['name'],
                    mime=doc_info['mime_type'],
                    key=f"download_{latest_doc_key}_{hash(doc_info['name'])}"
                )
                
                # Mark as handled
                st.session_state.doc_download_handled = True
    except Exception as e:
        print(f"Error handling document download: {e}")
    
    # Clear the container after handling
    # In a real implementation, we'd track which documents were already downloaded
    # For now, just clear after a short delay
    if st.session_state.doc_download_handled:
        # This empty container will replace the download button
        download_container.empty()

def main():
    # Debug print
    print("DEBUG: Starting main function")
    
    # Apply custom CSS
    st.markdown(CUSTOM_CSS, unsafe_allow_html=True)
    
    print("DEBUG: CSS applied")
    
    # Handle any document download requests
    handle_document_download_request()
    
    print("DEBUG: Document download handler registered")
    
    # Add banner
    st.markdown(
        """
        <div class="top-banner">
            <div class="banner-content">
                <div class="banner-text">KNOWTHEE.AI</div>
            </div>
        </div>
        <div class="top-spacer"></div>
        """,
        unsafe_allow_html=True
    )
    
    print("DEBUG: Banner added")
    
    # Initialize essential components if not already in session state
    print("DEBUG: About to initialize components")
    
    if 'employee_db' not in st.session_state:
        print("DEBUG: Initializing employee_db")
        st.session_state.employee_db = EmployeeDatabase()
        print("DEBUG: employee_db initialized")
    
    if 'query_processor' not in st.session_state:
        print("DEBUG: Initializing query_processor")
        st.session_state.query_processor = QueryProcessor()
        print("DEBUG: query_processor initialized")
        
    if 'vector_store' not in st.session_state:
        print("DEBUG: Initializing vector_store")
        st.session_state.vector_store = VectorStore()
        print("DEBUG: vector_store initialized")
        
        # IMPORTANT: Skip the initial loading of all employee profiles
        # This will happen on-demand when needed
        print("DEBUG: Skipping initial employee data loading to improve startup time")
    
    # Add a flag to track if employees are loaded
    if 'employees_loaded' not in st.session_state:
        st.session_state.employees_loaded = False
        
    print("DEBUG: Components initialized, setting up tabs")
    
    # Initialize active tab state if not already set
    if 'active_tab' not in st.session_state:
        st.session_state.active_tab = "Individual Profile"
    
    # Header container
    with st.container():
        st.markdown('<div class="header-bar">', unsafe_allow_html=True)
        st.markdown('<div class="header-title">AI Leadership Assessment</div>', unsafe_allow_html=True)
        st.markdown('<div class="header-subtitle">Upload documents to generate a personalized leadership profile</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
    
    print("DEBUG: Header added, creating tabs")
    
    # Add tabs for different functionality
    tab1, tab2 = st.tabs(["Individual Profile", "Employee Database"])
    
    print("DEBUG: Tabs created")
    
    # Set the active tab based on session state
    current_tab = st.session_state.active_tab
    
    print(f"DEBUG: Current tab: {current_tab}")
    
    # Individual Profile tab - existing functionality
    with tab1:
        print("DEBUG: Rendering tab1 (Individual Profile)")
        # Remove conditional rendering - always show content
        individual_profile_ui()
        # If this tab is active, update session state
        if current_tab == "Individual Profile":
            st.session_state.active_tab = "Individual Profile"
    
    print("DEBUG: Tab1 rendered")
    
    # Employee Database tab - new functionality
    with tab2:
        print("DEBUG: Rendering tab2 (Employee Database)")
        # Only now we load employee data if needed and if on Employee Database tab
        if not st.session_state.employees_loaded and current_tab == "Employee Database":
            with st.spinner("Loading employee database..."):
                load_employee_data()
                st.session_state.employees_loaded = True
                
        # Remove conditional rendering - always show content
        employee_database_ui()
        # If this tab is active, update session state
        if current_tab == "Employee Database":
            st.session_state.active_tab = "Employee Database"
    
    print("DEBUG: Tab2 rendered")
    
    # Add footer
    st.markdown(
        """
        <div class="footer">KNOW THYSELF & KNOW THY TEAM</div>
        """,
        unsafe_allow_html=True
    )
    
    print("DEBUG: Footer added, main function complete")

def load_employee_data():
    """Load employee data into vector store (separated to avoid blocking UI)"""
    print("DEBUG: Loading employee data")
    
    # Create a placeholder for progress
    progress_placeholder = st.empty()
    progress_bar = progress_placeholder.progress(0)
    
    # Reload employee profiles into the vector store
    employees = st.session_state.employee_db.get_all_employees()
    total_employees = len(employees)
    print(f"DEBUG: Found {total_employees} employees to load")
    
    # Skip if no employees
    if total_employees == 0:
        progress_placeholder.empty()
        return
    
    # Update progress message
    progress_placeholder.markdown(f"Loading {total_employees} employee profiles in batch mode...")
    progress_bar.progress(25)  # Show some progress immediately
    
    # Process in batch mode (much faster)
    try:
        # Get full employee data
        progress_placeholder.markdown(f"Retrieving employee data...")
        progress_bar.progress(50)
        
        employee_data_list = []
        for employee in employees:
            employee_id = employee['id']
            # Get the full employee profile
            employee_data = st.session_state.employee_db.get_employee(employee_id)
            if employee_data:
                employee_data_list.append(employee_data)
        
        # Use batch processing for significant speedup
        progress_placeholder.markdown(f"Processing {len(employee_data_list)} employees in batch mode...")
        progress_bar.progress(75)
        
        st.session_state.vector_store.batch_store_employee_profiles(employee_data_list)
        
        progress_bar.progress(100)
        progress_placeholder.markdown(f"✅ Successfully loaded {len(employee_data_list)} employee profiles!")
        time.sleep(1)  # Let the user see the success message
        
    except Exception as e:
        error_msg = f"Error loading employee data: {str(e)}"
        print(f"DEBUG: {error_msg}")
        progress_placeholder.error(error_msg)
        time.sleep(3)  # Let the user see the error
    
    # Clear progress indicator
    progress_placeholder.empty()
    print("DEBUG: Finished loading employee data")

def individual_profile_ui():
    """UI for the original individual profile functionality"""
    print("DEBUG: Starting individual_profile_ui")
    
    # Progress indicator
    current_step = 0
    if st.session_state.profile is not None:
        current_step = 2
    elif len(st.session_state.subject_docs) > 0:
        current_step = 1
    
    steps = ["Upload Documents", "Generate Profile", "View Profile"]
    st.markdown(f'<div class="progress-cue">Current Step: {steps[current_step]}</div>', unsafe_allow_html=True)
    
    print("DEBUG: Progress indicator rendered")
    
    # Document Upload Section
    with st.container():
        print("DEBUG: Rendering document upload section")
        
        st.markdown('<div class="section-title">Document Upload</div>', unsafe_allow_html=True)
        st.markdown('<div class="section-desc">Upload documents about the individual you want to assess.</div>', unsafe_allow_html=True)
        
        print("DEBUG: About to render file uploader")
        subject_files = st.file_uploader("Upload Leadership Documents", accept_multiple_files=True, type=['pdf', 'docx'], key='subject_upload')
        print("DEBUG: File uploader rendered")
        
        # Process documents if uploaded
        if subject_files:
            print("DEBUG: Subject files detected, processing")
            with st.spinner("Processing documents..."):
                # Clear existing docs if new ones are uploaded
                if subject_files != st.session_state.subject_docs:
                    # Cache each document for potential reference downloads
                    for doc in subject_files:
                        cache_document(doc)
                        
                    st.session_state.subject_docs = subject_files
                    st.session_state.profile = None
                    st.session_state.reference_docs = []
                
                # Show the list of uploaded documents
                if st.session_state.subject_docs:
                    st.markdown("### Uploaded Documents")
                    for doc in st.session_state.subject_docs:
                        st.write(f"- {doc.name}")
            print("DEBUG: Subject files processed")
    
    print("DEBUG: Document upload section rendered")
    
    # Profile Generation Section
    with st.container():
        print("DEBUG: Rendering profile generation section")
        
        st.markdown('<div class="section-title">Profile Generation</div>', unsafe_allow_html=True)
        st.markdown('<div class="section-desc">Generate a leadership profile based on uploaded documents.</div>', unsafe_allow_html=True)
        
        intent_option = st.selectbox(
            "What is your intent for this profile?",
            options=["Get an overall assessment", "Evaluate for a specific role", "Target development areas", "Other"],
            key="intent_select"
        )
        
        st.session_state.intent = intent_option
        
        if intent_option == "Other":
            st.session_state.intent_other = st.text_input("Please specify your intent:", key="intent_other_input")
        else:
            st.session_state.intent_other = ""
        
        generate_button = st.button("Generate Profile")
        
        if generate_button and st.session_state.subject_docs:
            print("DEBUG: Generate button clicked, has subject docs")
            with st.spinner("Generating leadership profile... This may take a minute."):
                # Process the documents and generate profile
                document_chunks = []
                metadata_list = []
                
                for doc in st.session_state.subject_docs:
                    # Create a temporary file
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{doc.name.split('.')[-1]}") as tmp:
                        tmp.write(doc.getvalue())
                        tmp_path = tmp.name
                    
                    # Process the document
                    doc_text, doc_metadata = document_processor.process_document(tmp_path)
                    
                    # Clean up the temporary file
                    os.unlink(tmp_path)
                    
                    # Append to our lists
                    document_chunks.append(doc_text)
                    metadata_list.append(doc_metadata)
                
                print("DEBUG: Documents processed, storing in vector DB")
                # Store in vector DB
                st.session_state.vector_store.store_documents(document_chunks, metadata_list)
                
                print("DEBUG: Saving as reference docs")
                # Save as reference docs
                st.session_state.reference_docs = document_chunks
                
                print("DEBUG: Generating profile")
                # Generate profile
                profile_json = profile_generator.generate_profile(document_chunks, metadata_list)
                st.session_state.profile = profile_json
                
                print("DEBUG: Profile generated, redirecting")
                # Redirect to show profile
                st.rerun()
    
    print("DEBUG: Profile generation section rendered")
    
    # Profile Display Section
    if st.session_state.profile is not None:
        print("DEBUG: Profile exists, displaying profile")
        display_profile(st.session_state.profile)
        
        # Question answering feature
        st.markdown('<div class="section-title">Ask Questions</div>', unsafe_allow_html=True)
        st.markdown('<div class="section-desc">Ask specific questions about this leadership profile.</div>', unsafe_allow_html=True)
        
        # Question input
        st.session_state.user_question = st.text_area("Ask a question about this leader:", key="question_input")
        
        # Submit button for the question
        if st.button("Submit Question"):
            if st.session_state.user_question:
                with st.spinner("Generating answer..."):
                    # Get relevant chunks from vector store
                    relevant_chunks = st.session_state.vector_store.get_relevant_chunks(st.session_state.user_question, n_results=5)
                    
                    # Combine with remaining docs to provide context
                    context_chunks = relevant_chunks
                    
                    # Generate answer
                    answer = profile_generator.answer_question(context_chunks, st.session_state.user_question)
                    st.session_state.question_answer = answer
        
        # Display the answer if available
        if st.session_state.question_answer:
            st.markdown('<div class="profile-section">', unsafe_allow_html=True)
            st.markdown("### Answer")
            
            # Process the answer to extract the main content and references
            answer_text = st.session_state.question_answer
            reference_section_start = answer_text.find("References")
            
            if reference_section_start > 0:
                # Split into main answer and references
                main_answer = answer_text[:reference_section_start].strip()
                references_text = answer_text[reference_section_start:].strip()
                
                # Display main answer content
                st.markdown(main_answer)
                
                # Process references to create clickable links
                st.markdown("### References")
                
                # Extract reference items
                reference_lines = []
                for line in references_text.split('\n'):
                    # Skip the "References" header line
                    if line.strip() and not line.strip().startswith('References'):
                        # Clean up bullet points
                        ref_item = line.strip().lstrip('•-* ')
                        if ref_item:
                            reference_lines.append(ref_item)
                
                # Create a dummy employee data structure with doc names from subject_docs
                if st.session_state.subject_docs:
                    dummy_employee = {
                        "id": "individual",
                        "metadata": {
                            "document_names": [doc.name for doc in st.session_state.subject_docs]
                        }
                    }
                    
                    # Create clickable references
                    if reference_lines:
                        reference_html = create_reference_links(dummy_employee, reference_lines)
                        st.markdown(reference_html, unsafe_allow_html=True)
            else:
                # No references section found, display the whole answer
                st.markdown(answer_text)
            
            st.markdown('</div>', unsafe_allow_html=True)
        
        # Download buttons
        col1, col2 = st.columns(2)
        
        with col1:
            # PDF download - try to generate but handle errors silently
            try:
                pdf_data = create_pdf(st.session_state.profile, st.session_state.question_answer)
                st.download_button(
                    label="Download PDF Report",
                    data=pdf_data,
                    file_name="leadership_profile.pdf",
                    mime="application/pdf",
                )
            except Exception as e:
                # Just log the error but don't display to user
                print(f"DEBUG: Unable to generate PDF: {str(e)}")
        
        with col2:
            # Create and offer PowerPoint download
            try:
                pptx_bytes = generate_pptx_from_profile(st.session_state.profile)
                st.download_button(
                    label="Download Leadership Profile",
                    data=pptx_bytes,
                    file_name="leadership_profile.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as e:
                # Just log the error but don't display to user
                print(f"DEBUG: Unable to generate PowerPoint: {str(e)}")
    
    print("DEBUG: individual_profile_ui completed")

def employee_database_ui():
    """UI for the employee database functionality"""
    st.markdown('<div class="section-title">Employee Database</div>', unsafe_allow_html=True)
    st.markdown('<div class="section-desc">Manage and search your employee profiles database.</div>', unsafe_allow_html=True)
    
    # Access the vector store from session state
    vector_store = st.session_state.vector_store
    
    # Check if data is loaded
    if not st.session_state.get('employees_loaded', False):
        load_button = st.button("Load Employee Database")
        if load_button:
            with st.spinner("Loading employee database..."):
                load_employee_data()
                st.session_state.employees_loaded = True
                # Force refresh to show the database tabs
                st.rerun()
        
        st.info("Please load the employee database to access the search and management functionality.")
        return  # Exit early if data not loaded
    
    # Create tabs for different database operations
    db_tab1, db_tab2, db_tab3, db_tab4 = st.tabs(["Search Employees", "Intelligent Queries", "Add to Database", "Manage Database"])
    
    # Tab 1: Search Employees
    with db_tab1:
        st.markdown("### Natural Language Search")
        st.markdown("Search for employees using natural language queries like 'extroverted, creative engineers'")
        st.markdown("After finding employees, you can ask specific questions about them in their profile cards.")
        
        # Search input
        search_query = st.text_input("Enter your search query:", key="employee_search_query")
        
        # Initialize search_results in session state if not present
        if 'search_results' not in st.session_state:
            st.session_state.search_results = None
            
        # Initialize employee_answers in session state if not present
        if 'employee_answers' not in st.session_state:
            st.session_state.employee_answers = {}
            
        if st.button("Search"):
            if search_query:
                with st.spinner("Searching employees..."):
                    # Parse the query
                    parsed_query = st.session_state.query_processor.parse_query(search_query)
                    
                    # Convert to filters
                    filters = st.session_state.query_processor.convert_to_filters(parsed_query)
                    
                    # Execute search
                    results = vector_store.search_employees(search_query, filters)
                    
                    # Process and display results
                    if results:
                        processed_results = st.session_state.query_processor.process_search_results(
                            results, search_query, parsed_query
                        )
                        
                        # Store results in session
                        st.session_state.search_results = processed_results
                    
        # Display search results if available in session state
        if st.session_state.search_results is not None:
            # Display explanation
            st.markdown("### Search Results")
            st.markdown(f"*{st.session_state.search_results['explanation']}*")
            
            # Display employee cards with expanded functionality
            if st.session_state.search_results["count"] > 0:
                employees = st.session_state.search_results["employees"]
                
                # For each employee, create an expanded card with all functionality
                for i, employee in enumerate(employees):
                    employee_id = employee['id']
                    
                    # Get full employee data
                    employee_data = st.session_state.employee_db.get_employee(employee_id)
                    if not employee_data:
                        continue
                    
                    with st.expander(f"{employee['name']} - {', '.join(employee['traits'])}", expanded=True):
                        # Basic info section
                        st.markdown(f"### {employee['name']}")
                        
                        # Display metadata
                        metadata = employee_data['metadata']
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.markdown("**Department:** " + metadata.get('department', 'Not specified'))
                            
                            if 'traits' in metadata and metadata['traits']:
                                if isinstance(metadata['traits'], list):
                                    traits = ", ".join(metadata['traits'])
                                else:
                                    traits = metadata['traits']
                                st.markdown(f"**Traits:** {traits}")
                        
                        with col2:
                            if 'roles' in metadata and metadata['roles']:
                                if isinstance(metadata['roles'], list):
                                    roles = ", ".join(metadata['roles'])
                                else:
                                    roles = metadata['roles']
                                st.markdown(f"**Roles:** {roles}")
                            
                            if 'leadership_style' in metadata and metadata['leadership_style']:
                                if isinstance(metadata['leadership_style'], list):
                                    style = ", ".join(metadata['leadership_style'])
                                else:
                                    style = metadata['leadership_style']
                                st.markdown(f"**Leadership Style:** {style}")
                        
                        # Display document downloads section
                        st.markdown("### Source Documents")
                        
                        # Extract source files section
                        profile_json = employee_data['profile']
                        try:
                            # Handle both string and already parsed JSON
                            if isinstance(profile_json, str):
                                profile_data = json.loads(profile_json)
                            else:
                                profile_data = profile_json
                        except Exception as e:
                            print(f"Error parsing profile JSON: {str(e)}")
                            profile_data = []
                        
                        # Extract unique source files from all profile sections
                        source_files = set()
                        for section in profile_data:
                            sources = section.get('sources', '')
                            if sources:
                                # Clean up sources text
                                clean_sources = clean_source_text(sources)
                                # Add each source to the set
                                for source in clean_sources.split(','):
                                    source = source.strip()
                                    if source:
                                        source_files.add(source)
                        
                        # Display source files if any
                        if source_files:
                            st.markdown("#### Documents Used to Generate This Profile")
                            for source in sorted(source_files):
                                st.markdown(f"- {source}")
                        
                        # Display uploaded documents if available
                        if 'document_names' in employee_data['metadata'] and employee_data['metadata']['document_names']:
                            st.markdown("#### Original Documents")
                            for doc_name in employee_data['metadata']['document_names']:
                                st.markdown(f"- {doc_name}")
                        
                        # Add download buttons for profile
                        st.markdown("### Download Profile")
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            # PDF download - try to create but silently handle errors
                            try:
                                pdf_data = create_pdf(profile_json)
                                st.download_button(
                                    label="Download PDF Report",
                                    data=pdf_data,
                                    file_name=f"{employee_data['name']}_profile.pdf",
                                    mime="application/pdf",
                                    key=f"pdf_{employee_id}"
                                )
                            except Exception as e:
                                # Just log the error but don't display to user
                                print(f"Unable to generate PDF: {str(e)}")
                        
                        with col2:
                            # Create and offer PowerPoint download
                            try:
                                pptx_bytes = generate_pptx_from_profile(profile_json)
                                st.download_button(
                                    label="Download Leadership Profile",
                                    data=pptx_bytes,
                                    file_name=f"{employee_data['name']}_profile.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key=f"pptx_{employee_id}"
                                )
                            except Exception as e:
                                # Just log the error but don't display to user
                                print(f"Unable to generate PowerPoint: {str(e)}")
                        
                        # Add question answering section
                        st.markdown('<div class="section-title">Ask Questions</div>', unsafe_allow_html=True)
                        st.markdown(f'Ask specific questions about {employee_data["name"]}')
                        
                        question_key = f"question_{employee_id}"
                        st.markdown("Your question:", unsafe_allow_html=True)
                        employee_question = st.text_area("", 
                            placeholder=f"Example: What are {employee_data['name']}'s key strengths?",
                            key=question_key)
                        
                        # Initialize this employee's answers if not present
                        if employee_id not in st.session_state.employee_answers:
                            st.session_state.employee_answers[employee_id] = None
                            
                        if st.button("Submit Question", key=f"submit_{employee_id}"):
                            if employee_question:
                                with st.spinner(f"Generating answer..."):
                                    try:
                                        # Get relevant chunks from vector store for more detailed citations
                                        relevant_chunks = st.session_state.vector_store.get_relevant_chunks(
                                            employee_question, 
                                            n_results=5, 
                                            employee_id=employee_id
                                        )
                                        
                                        # If no relevant chunks found, fall back to profile sections
                                        if not relevant_chunks:
                                            context_chunks = []
                                            for section in profile_data:
                                                section_text = f"{section.get('section', '')}: {section.get('content', '')}"
                                                context_chunks.append(section_text)
                                        else:
                                            context_chunks = relevant_chunks
                                        
                                        # Generate answer
                                        answer = profile_generator.answer_question(context_chunks, employee_question)
                                        
                                        # Store answer in the centralized employee_answers dict
                                        st.session_state.employee_answers[employee_id] = answer
                                    except Exception as e:
                                        print(f"Error generating answer: {str(e)}")
                                        st.session_state.employee_answers[employee_id] = f"Sorry, I couldn't generate an answer due to an error. Please try a different question."
                        
                                            # Display the answer if available for this employee
                    if st.session_state.employee_answers[employee_id] is not None:
                        st.markdown('<div class="profile-section">', unsafe_allow_html=True)
                        st.markdown("### Answer")
                        
                        # Process the answer to extract the main content and references
                        answer_text = st.session_state.employee_answers[employee_id]
                        reference_section_start = answer_text.find("References")
                        
                        if reference_section_start > 0:
                            # Split into main answer and references
                            main_answer = answer_text[:reference_section_start].strip()
                            references_text = answer_text[reference_section_start:].strip()
                            
                            # Display main answer content
                            st.markdown(main_answer)
                            
                            # Process references to create clickable links
                            st.markdown("### References")
                            
                            # Extract reference items
                            reference_lines = []
                            for line in references_text.split('\n'):
                                # Skip the "References" header line
                                if line.strip() and not line.strip().startswith('References'):
                                    # Clean up bullet points
                                    ref_item = line.strip().lstrip('•-* ')
                                    if ref_item:
                                        reference_lines.append(ref_item)
                            
                            # Create clickable references
                            if reference_lines:
                                reference_html = create_reference_links(employee_data, reference_lines)
                                st.markdown(reference_html, unsafe_allow_html=True)
                        else:
                            # No references section found, display the whole answer
                            st.markdown(answer_text)
                        
                        st.markdown('</div>', unsafe_allow_html=True)
            else:
                st.markdown("No matching employees found.")
        else:
            if search_query:  # Only show this message if user has searched
                st.markdown("No results found. Try a different query or add more employees to the database.")
    
    # Tab 2: Intelligent Queries (Enhanced RAG System)
    with db_tab2:
        st.markdown("### 🧠 Intelligent HR Analytics")
        st.markdown("Ask complex questions with intelligent conversation context management.")
        
        # === NEW: Conversation Settings Panel === (HIDDEN)
        # with st.expander("⚙️ Conversation Settings", expanded=False):
        #     col1, col2, col3 = st.columns(3)
        #     
        #     with col1:
        #         memory_mode = st.selectbox(
        #             "Memory Mode:",
        #             options=["adaptive", "short", "medium", "long"],
        #             index=0,
        #             help="Controls how much conversation history to maintain"
        #         )
        #     
        #     with col2:
        #         focus_mode = st.selectbox(
        #             "Employee Focus:",
        #             options=["adaptive", "narrow", "broad"],
        #             index=0,
        #             help="Controls how many employees to consider in analysis"
        #         )
        #     
        #     with col3:
        #         enable_context = st.checkbox(
        #             "Context Tracking",
        #             value=True,
        #             help="Enable intelligent conversation context tracking"
        #         )
        #     
        #     # Update settings if changed
        #     new_settings = {
        #         "max_conversation_memory": memory_mode,
        #         "employee_focus_mode": focus_mode,
        #         "enable_context_tracking": enable_context
        #     }
        #     rag_system.update_conversation_settings(new_settings)
        
        # === Enhanced Conversation Status Display ===
        conversation_insights = rag_system.get_conversation_insights()
        
        # Commenting out the status display - HIDDEN BY USER REQUEST
        # if conversation_insights.get("status") != "no_conversation":
        #     # Main conversation status bar
        #     col1, col2, col3, col4 = st.columns([2, 1, 1, 1])
        #     
        #     with col1:
        #         memory_usage = conversation_insights["memory_usage"]
        #         memory_color = "🟢" if memory_usage < 50 else "🟡" if memory_usage < 80 else "🔴"
        #         st.metric(
        #             "Conversation Status",
        #             f"{conversation_insights['conversation_length']} exchanges",
        #             f"{memory_color} {memory_usage:.0f}% memory"
        #         )
        #     
        #     with col2:
        #         st.metric(
        #             "Context Employees",
        #             conversation_insights["current_context_employees"],
        #             f"Theme: {conversation_insights['conversation_theme']}"
        #         )
        #     
        #     with col3:
        #         tokens_used = conversation_insights["total_tokens_used"]
        #         st.metric(
        #             "Tokens Used",
        #             f"{tokens_used:,}",
        #             "Total conversation"
        #         )
        #     
        #     with col4:
        #         if st.button("🗑️ Clear", help="Start fresh conversation", key="clear_conv"):
        #             rag_system.clear_conversation_history()
        #             st.success("Conversation cleared!")
        #             st.rerun()
        
        # === Enhanced Conversation History ===
        if conversation_insights.get("status") != "no_conversation":
            with st.expander(f"💬 Conversation History ({conversation_insights['conversation_length']} exchanges)", expanded=False):
                if hasattr(rag_system, 'conversation_history') and rag_system.conversation_history:
                    for i, entry in enumerate(rag_system.conversation_history[-5:], 1):  # Show last 5
                        col_q, col_a = st.columns([1, 2])
                        
                        with col_q:
                            st.markdown(f"**Q{i}:** {entry['original_query']}")
                            if entry.get('resolved_query') != entry['original_query']:
                                st.markdown(f"*→ {entry['resolved_query']}*")
                            if entry.get('context_employees'):
                                st.markdown(f"*Context: {', '.join(entry['context_employees'][:2])}*")
                        
                        with col_a:
                            response_preview = entry['response'][:200] + "..." if len(entry['response']) > 200 else entry['response']
                            st.markdown(f"**A{i}:** {response_preview}")
                            
                            # Show query metadata
                            query_meta = f"Type: {entry.get('query_type', 'general')} • Tokens: {entry.get('tokens_used', 0)}"
                            st.caption(query_meta)
                        
                        st.markdown("---")
            
            # === Context Employees Display ===
            if conversation_insights.get("top_employees"):
                with st.expander("👥 Employee Context", expanded=False):
                    st.markdown("**Most discussed employees:**")
                    for emp_name, frequency in conversation_insights["top_employees"]:
                        st.markdown(f"• **{emp_name}** ({frequency} mentions)")
                    
                    # Show current context with relevance scores
                    if hasattr(rag_system, 'context_employees') and rag_system.context_employees:
                        st.markdown("**Current context employees:**")
                        for emp in rag_system.context_employees[:8]:  # Show top 8
                            if isinstance(emp, dict):
                                score_bar = "🟩" * int(emp["relevance_score"] * 5) + "⬜" * (5 - int(emp["relevance_score"] * 5))
                                st.markdown(f"• {emp['name']} {score_bar} ({emp['relevance_score']:.1f})")
        else:
            # First time user guidance
            st.info("💡 Start a conversation! Your questions will build context intelligently.")
        
        # === Example Queries with Dynamic Suggestions ===
        with st.expander("💡 Example Queries & Smart Follow-ups", expanded=not conversation_insights.get("status") != "no_conversation"):
            if conversation_insights.get("status") == "no_conversation":
                st.markdown("**Start with these analysis types:**")
                examples = {
                    "🎯 Succession Planning": [
                        "Who are the top 3 candidates for senior engineering leadership?",
                        "Which employees show the highest leadership potential in marketing?"
                    ],
                    "👥 Team Analysis": [
                        "What are the strengths and gaps in our product team?",
                        "Who would work well together on an innovation project?"
                    ],
                    "⚠️ Risk Assessment": [
                        "Which high-performers might be at risk of leaving?",
                        "What are common derailers among senior leaders?"
                    ],
                    "📈 Development Planning": [
                        "Which employees need communication skills development?",
                        "Who has potential for international assignments?"
                    ]
                }
            else:
                st.markdown("**Smart follow-up suggestions based on your conversation:**")
                theme = conversation_insights.get("conversation_theme", "general")
                context_employees = [emp["name"] if isinstance(emp, dict) else emp 
                                   for emp in rag_system.context_employees[:3]]
                
                if context_employees and len(context_employees) >= 2:
                    examples = {
                        "🔄 Compare Current Context": [
                            f"Between {context_employees[0]} and {context_employees[1]}, who is more creative?",
                            "What about their leadership styles?",
                            "Which among them would be better for a client-facing role?"
                        ],
                        "📊 Expand Analysis": [
                            "Who else in the organization has similar strengths?",
                            "What are the development needs across this group?",
                            "How do they compare to industry benchmarks?"
                        ]
                    }
                else:
                    examples = {
                        "🎯 Continue Theme": [
                            f"Who else should be considered for {theme.replace('_', ' ')}?",
                            "What patterns do you see across these profiles?",
                            "What are the key risks we should be aware of?"
                        ]
                    }
            
            for category, queries in examples.items():
                st.markdown(f"**{category}:**")
                for query in queries:
                    col1, col2 = st.columns([4, 1])
                    with col1:
                        st.markdown(f"• {query}")
                    with col2:
                        if st.button("Use", key=f"use_{hash(query)}", help="Use this query"):
                            st.session_state.intelligent_query_input = query
                            st.rerun()
        
        # === Smart Query Input with Context Hints ===
        st.markdown("**Enter your HR analytics question:**")
        
        # Show context hints
        if hasattr(rag_system, 'context_employees') and rag_system.context_employees:
            context_names = [emp["name"] if isinstance(emp, dict) else emp for emp in rag_system.context_employees[:3]]
            hint_text = f"💡 *You can refer to: {', '.join(context_names)}"
            if len(rag_system.context_employees) > 3:
                hint_text += f" and {len(rag_system.context_employees) - 3} others"
            hint_text += " using 'them', 'between them', 'those employees', etc.*"
            st.markdown(hint_text)
        
        # Query input with preserved state
        if 'intelligent_query_input' not in st.session_state:
            st.session_state.intelligent_query_input = ""
        
        intelligent_query = st.text_area(
            "",
            value=st.session_state.intelligent_query_input,
            height=100,
            placeholder="Example: Analyze leadership pipeline in engineering and identify succession opportunities.",
            key="intelligent_query_textarea",
            label_visibility="collapsed"
        )
        
        # Update session state when text changes
        if intelligent_query != st.session_state.intelligent_query_input:
            st.session_state.intelligent_query_input = intelligent_query
        
        # === Enhanced Query Configuration === (REMOVED)
        # col1, col2 = st.columns([2, 1])
        # 
        # with col1:
        #     query_type = st.selectbox(
        #         "Analysis Type:",
        #         options=[
        #             "General Analysis",
        #             "Succession Planning", 
        #             "Team Composition",
        #             "Risk Assessment",
        #             "Development Planning",
        #             "Cross-Employee Comparison",
        #             "Department Analysis"
        #         ],
        #         key="query_type_select",
        #         help="Helps optimize analysis scope and employee limits"
        #     )
        # 
        # with col2:
        #     # Show estimated employee scope based on query type and settings
        #     query_type_key = query_type.lower().replace(" ", "_")
        #     if query_type_key in rag_system.employee_limits:
        #         limits = rag_system._get_employee_limit_for_query(query_type_key, "multiple_employees")
        #         est_employees = f"{limits['priority']}-{limits['max']} employees"
        #     else:
        #         est_employees = "5-10 employees"
        #     
        #     st.metric("Estimated Scope", est_employees, "Will analyze")
        
        # === Enhanced Analysis Button and Processing ===
        if 'intelligent_query_results' not in st.session_state:
            st.session_state.intelligent_query_results = None
        
        if st.button("🔍 Analyze", key="intelligent_query_submit", type="primary"):
            if intelligent_query:
                with st.spinner("🧠 Processing with intelligent context management..."):
                    try:
                        # Clear the input after submission
                        st.session_state.intelligent_query_input = ""
                        
                        # Use the RAG system for intelligent analysis with automatic query type detection
                        result = rag_system.process_complex_query(
                            query=intelligent_query,
                            context_type="general_analysis"  # Let the system auto-detect the appropriate type
                        )
                        
                        st.session_state.intelligent_query_results = result
                        
                    except Exception as e:
                        st.error(f"Analysis failed: {str(e)}")
                        print(f"RAG query error: {e}")
        
        # === Enhanced Results Display === (HIDDEN BY USER REQUEST)
        # if st.session_state.intelligent_query_results:
        #     result = st.session_state.intelligent_query_results
        #     
        #     st.markdown("---")
        #     st.markdown("### 📊 Analysis Results")
        #     
        #     # === Query Resolution Display ===
        #     if result.get('resolved_query') != result.get('query'):
        #         with st.expander("🔄 Query Resolution", expanded=True):
        #             col1, col2 = st.columns([1, 1])
        #             with col1:
        #                 st.markdown(f"**Original:** {result.get('query', '')}")
        #             with col2:
        #                 st.markdown(f"**Understood as:** {result.get('resolved_query', '')}")
        #             
        #             if result.get('context_employees'):
        #                 st.markdown(f"**Context employees:** {', '.join(result['context_employees'])}")
        #             
        #             # Show conversation status update
        #             conv_status = result.get('conversation_status', {})
        #             if conv_status:
        #                 st.markdown(f"**Conversation:** {conv_status.get('conversation_length', 0)} exchanges, "
        #                           f"{conv_status.get('context_employees_count', 0)} context employees")
        #     
        #     # === Analysis Metadata ===
        #     with st.expander("🔍 Analysis Details", expanded=False):
        #         analysis = result.get('analysis', {})
        #         col1, col2, col3 = st.columns(3)
        #         
        #         with col1:
        #             st.markdown(f"**Query Type:** {analysis.get('query_type', 'N/A')}")
        #             st.markdown(f"**Scope:** {analysis.get('scope', 'N/A')}")
        #             st.markdown(f"**Analysis Depth:** {analysis.get('analysis_depth', 'N/A')}")
        #         
        #         with col2:
        #             st.markdown(f"**Data Sources:** {', '.join(analysis.get('required_data', ['general']))}")
        #             st.markdown(f"**Context Sources:** {result.get('context_sources', 0)}")
        #             
        #             # Show employee limits used
        #             emp_limits = result.get('employee_limits', {})
        #             st.markdown(f"**Employee Limits:** {emp_limits.get('priority', 'N/A')}/{emp_limits.get('max', 'N/A')}")
        #         
        #         with col3:
        #             # Show conversation statistics
        #             conv_status = result.get('conversation_status', {})
        #             if conv_status:
        #                 memory_status = conv_status.get('memory_status', {})
        #                 st.markdown(f"**Memory Usage:** {memory_status.get('usage_percentage', 0):.0f}%")
        #                 st.markdown(f"**Token Limit:** {memory_status.get('token_limit', 0):,}")
        #                 
        #                 settings = conv_status.get('settings', {})
        #                 st.markdown(f"**Mode:** {settings.get('memory_mode', 'adaptive')}/{settings.get('focus_mode', 'adaptive')}")
        #     
        #     # === Main Response with Enhanced Formatting ===
        #     response = result.get('response', '')
        #     if response:
        #         st.markdown('<div class="profile-section">', unsafe_allow_html=True)
        #         
        #         # Process response for better formatting
        #         if "References" in response:
        #             main_response, references = response.split("References", 1)
        #             st.markdown(main_response.strip())
        #             
        #             with st.expander("📚 Sources & Citations", expanded=False):
        #                 st.markdown("**References**" + references)
        #         else:
        #             st.markdown(response)
        #         
        #         st.markdown('</div>', unsafe_allow_html=True)
        #         
        #         # === Enhanced Action Buttons ===
        #         col1, col2, col3 = st.columns([1, 1, 1])
        #         
        #         with col1:
        #             if st.button("📥 Download Analysis", key="download_analysis"):
        #                 timestamp = time.strftime('%Y%m%d_%H%M%S')
        #                 report_content = f"""# HR Analytics Report
        # 
        # ## Query
        # {intelligent_query}
        # 
        # ## Analysis Type
        # {query_type}
        # 
        # ## Employee Context
        # {', '.join(result.get('context_employees', [])) if result.get('context_employees') else 'None'}
        # 
        # ## Results
        # {response}
        # 
        # ## Analysis Metadata
        # - Query Type: {analysis.get('query_type', 'N/A')}
        # - Scope: {analysis.get('scope', 'N/A')}
        # - Sources: {result.get('context_sources', 0)}
        # - Employee Limits: {emp_limits.get('priority', 'N/A')}/{emp_limits.get('max', 'N/A')}
        # 
        # ## Generated
        # {time.strftime('%Y-%m-%d %H:%M:%S')}
        # """
        #                 st.download_button(
        #                     label="Download Report",
        #                     data=report_content,
        #                     file_name=f"hr_analysis_{timestamp}.txt",
        #                     mime="text/plain",
        #                     key="download_btn"
        #                 )
        #         
        #         with col2:
        #             if st.button("🔄 Ask Follow-up", key="ask_followup"):
        #                 st.session_state.intelligent_query_input = "What about their "
        #                 st.rerun()
        #         
        #         with col3:
        #             if st.button("📊 Deep Dive", key="deep_dive"):
        #                 context_emps = result.get('context_employees', [])
        #                 if context_emps:
        #                     followup = f"Provide a detailed analysis of {', '.join(context_emps[:2])} including specific assessment scores and development recommendations"
        #                     st.session_state.intelligent_query_input = followup
        #                     st.rerun()
        #     else:
        #         st.warning("No response generated. Please try a different query.")
        
        # Simple response display (replaces the complex analysis results section)
        if st.session_state.intelligent_query_results:
            result = st.session_state.intelligent_query_results
            response = result.get('response', '')
            if response:
                st.markdown("---")
                st.markdown('<div class="profile-section">', unsafe_allow_html=True)
                
                # Process response for better formatting
                if "References" in response:
                    main_response, references = response.split("References", 1)
                    st.markdown(main_response.strip())
                    
                    with st.expander("📚 Sources & Citations", expanded=False):
                        st.markdown("**References**" + references)
                else:
                    st.markdown(response)
                
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                st.warning("No response generated. Please try a different query.")
    
    # Tab 3: Add to Database
    with db_tab3:
        st.markdown("### Add Employee Profile")
        st.markdown("Generate and add a new employee profile to the database.")
        
        # Employee name input
        employee_name = st.text_input("Employee Name:", key="new_employee_name")
        
        # Document upload
        new_employee_files = st.file_uploader(
            "Upload Employee Documents", 
            accept_multiple_files=True, 
            type=['pdf', 'docx'], 
            key='new_employee_upload'
        )
        
        # Department selection
        department = st.selectbox(
            "Department:", 
            options=["Engineering", "Marketing", "Sales", "Product", "Design", "HR", "Finance", "Operations", "Other"],
            key="department_select"
        )
        
        if department == "Other":
            department = st.text_input("Specify department:", key="department_other")
        
        # Add to database button
        if st.button("Generate Profile & Add to Database"):
            if employee_name and new_employee_files:
                with st.spinner("Processing documents and generating profile..."):
                    # Process documents
                    document_chunks = []
                    metadata_list = []
                    
                    # Track document names
                    document_names = []
                    
                    for doc in new_employee_files:
                        # Track document name
                        document_names.append(doc.name)
                        
                        # Cache the document for reference downloads
                        cache_document(doc, employee_id="new")  # We'll update this ID after creation
                        
                        # Create a temporary file
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{doc.name.split('.')[-1]}") as tmp:
                            tmp.write(doc.getvalue())
                            tmp_path = tmp.name
                            
                            # Process the document
                            doc_text, doc_metadata = document_processor.process_document(tmp_path)
                            
                            # Clean up the temporary file
                            os.unlink(tmp_path)
                            
                            # Append to our lists
                            document_chunks.append(doc_text)
                            metadata_list.append(doc_metadata)
                    
                    # Generate profile
                    profile_json = profile_generator.generate_profile(document_chunks, metadata_list)
                    
                    # Add metadata
                    additional_metadata = {
                        "department": department,
                        "document_names": document_names
                    }
                    
                    # Add to employee database
                    employee_id = st.session_state.employee_db.add_employee(
                        name=employee_name,
                        profile_data=profile_json,
                        metadata=additional_metadata
                    )
                    
                    # Add to vector store for searching
                    profile_data = json.loads(profile_json)
                    employee_metadata = st.session_state.employee_db.get_employee(employee_id)["metadata"]
                    
                    # Store the profile sections for profile display
                    st.session_state.vector_store.store_employee_profile(
                        employee_id=employee_id,
                        profile_sections=profile_data,
                        metadata=employee_metadata
                    )
                    
                    # Also store the raw document chunks for detailed question answering
                    # This allows more detailed citations with actual assessment scores
                    st.session_state.vector_store.store_employee_documents(
                        employee_id=employee_id,
                        documents=document_chunks,
                        metadata=employee_metadata
                    )
                    
                    st.success(f"Successfully added {employee_name} to the employee database!")
                    
                    # Show profile summary
                    st.markdown("### Profile Summary")
                    
                    # Display the first section (Profile Summary)
                    if profile_data and len(profile_data) > 0:
                        summary_section = next((s for s in profile_data if s.get("section") == "Profile Summary"), None)
                        if summary_section:
                            st.markdown(summary_section.get("content", ""))
    
    # Tab 4: Manage Database
    with db_tab4:
        st.markdown("### Employee Database Management")
        
        # Get all employees
        employees = st.session_state.employee_db.get_all_employees()
        
        if not employees:
            st.markdown("No employees in the database yet. Add some using the 'Add to Database' tab.")
        else:
            # Create a DataFrame for display
            employee_data = []
            for emp in employees:
                # Extract traits as a comma-separated string
                traits = ", ".join(emp["metadata"].get("traits", []))
                
                # Format the added date
                added_date = emp["added_date"]
                if isinstance(added_date, str):
                    # Just take the date part if it's in ISO format
                    added_date = added_date.split("T")[0]
                
                employee_data.append({
                    "ID": emp["id"],
                    "Name": emp["name"],
                    "Department": emp["metadata"].get("department", ""),
                    "Traits": traits,
                    "Added Date": added_date
                })
            
            df = pd.DataFrame(employee_data)
            
            # Display as a table
            st.dataframe(df)
            
            # Allow deletion of employees
            st.markdown("### Delete Employee")
            
            # Create a dropdown with employee names
            employee_to_delete = st.selectbox(
                "Select employee to delete:",
                options=[f"{emp['name']} ({emp['id']})" for emp in employees],
                key="delete_employee_select"
            )
            
            if st.button("Delete Selected Employee"):
                # Extract the ID from the selection
                employee_id = employee_to_delete.split("(")[-1].split(")")[0]
                
                # Delete from database and vector store
                if st.session_state.employee_db.delete_employee(employee_id):
                    st.session_state.vector_store.delete_employee_profile(employee_id)
                    st.success(f"Successfully deleted employee from the database.")
                    st.rerun()
                else:
                    st.error("Failed to delete employee. Please try again.")

if __name__ == "__main__":
    main()
